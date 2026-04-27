"""
PPTX Generator — Creates PowerPoint presentations using python-pptx.
"""
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def generate_ppt_file(ppt_data: dict) -> io.BytesIO:
    """Generate a PowerPoint file from structured slide data."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_text = ppt_data.get("title", "Presentation")
    slides_data = ppt_data.get("slides", [])

    # ── Title Slide ──
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0f, 0x34, 0x60)

    # Title text box
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(11.333)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # ── Content Slides ──
    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # Heading bar
        left = Inches(0)
        top = Inches(0)
        width = Inches(13.333)
        height = Inches(1.2)
        shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x0f, 0x34, 0x60)
        shape.line.fill.background()

        # Heading text
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_info.get("heading", "")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.LEFT
        tf.margin_left = Inches(0.5)
        tf.margin_top = Inches(0.25)

        # Bullet points
        points = slide_info.get("points", [])
        left = Inches(0.8)
        top = Inches(1.6)
        width = Inches(11.5)
        height = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"•  {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0x2d, 0x2d, 0x2d)
            p.space_after = Pt(12)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
